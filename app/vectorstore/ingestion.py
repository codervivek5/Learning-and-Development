import os
import uuid
from typing import List
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from app.vectorstore.chroma_client import get_chroma_client
from app.vectorstore.embeddings import get_embeddings
from app.core.logging import get_logger

logger = get_logger(__name__)


def parse_pdf(file_path: str) -> str:
    """Extract text from a PDF file using PyMuPDF."""
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text


def parse_docx(file_path: str) -> str:
    """Extract text from a Word document (.docx)."""
    doc = DocxDocument(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])


def parse_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint presentation (.pptx)."""
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
    return "\n".join(text_runs)


def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks."""
    chunks = []
    if not text:
        return chunks

    start = 0
    text_len = len(text)

    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunks.append(text[start:end])
        if end == text_len:
            break
        start += chunk_size - chunk_overlap
    return chunks


async def ingest_document(
    project_id: uuid.UUID,
    file_path: str,
    document_id: uuid.UUID,
) -> bool:
    """Ingest a document: parse, chunk, embed, and store in ChromaDB."""
    try:
        if not os.path.exists(file_path):
            logger.error("File does not exist for ingestion", path=file_path)
            return False

        # Parse text based on extension
        _, ext = os.path.splitext(file_path.lower())
        logger.info("Parsing document", path=file_path, extension=ext)

        if ext == ".pdf":
            text = parse_pdf(file_path)
        elif ext == ".docx":
            text = parse_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            text = parse_pptx(file_path)
        else:
            # Fallback to plain text
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        chunks = chunk_text(text)
        if not chunks:
            logger.warning("No text extracted from document", path=file_path)
            return False

        logger.info("Generated text chunks", count=len(chunks))

        # Generate embeddings
        embeddings = await get_embeddings(chunks)

        # Get Chroma client
        client = get_chroma_client()
        collection_name = f"project_{str(project_id).replace('-', '_')}"

        # Get or create the collection for this project
        collection = client.get_or_create_collection(name=collection_name)

        # Prepare batch insert data
        ids = [f"{str(document_id)}_{i}" for i in range(len(chunks))]
        metadatas = [{"document_id": str(document_id), "project_id": str(project_id)} for _ in chunks]

        # Insert chunks
        collection.add(
            embeddings=embeddings,
            documents=chunks,
            ids=ids,
            metadatas=metadatas,
        )

        logger.info("Successfully ingested document", document_id=document_id, project_id=project_id)
        return True

    except Exception as e:
        logger.error("Ingestion failed", document_id=document_id, project_id=project_id, error=str(e))
        return False
