from datetime import datetime
from pathlib import Path
from typing import Optional

from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from sqlalchemy.ext.asyncio import AsyncSession

from app.services.project_service import ProjectService
from app.services.training_content_service import TrainingContentService

EXPORT_DIR = Path("./storage/exports")
EXPORT_DIR.mkdir(parents=True, exist_ok=True)


class ExportService:
    @staticmethod
    async def export_pdf(
        db: AsyncSession,
        project_id: int,
    ) -> str:
        review_record = await TrainingContentService.get_latest_phase_output(
            db=db,
            project_id=project_id,
            phase="review",
        )
        if not review_record:
            raise RuntimeError("Review phase output is required for export.")

        project = await ProjectService.get_project(db, project_id)
        if not project:
            raise RuntimeError("Project not found.")

        filename = f"project_{project_id}_review_export_{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.pdf"
        file_path = EXPORT_DIR / filename

        text_lines = [
            f"Project: {project.title}",
            f"Review Score: {review_record.content.get('score', 'N/A')}",
            "\nStrengths:",
        ]
        text_lines.extend([f"- {item}" for item in review_record.content.get('strengths', [])])
        text_lines.append("\nImprovements:")
        text_lines.extend([f"- {item}" for item in review_record.content.get('improvements', [])])
        text_lines.append(f"\nApproved: {review_record.content.get('approved', False)}")

        c = canvas.Canvas(str(file_path), pagesize=letter)
        width, height = letter
        y = height - 40

        for line in text_lines:
            c.drawString(40, y, str(line))
            y -= 18
            if y < 40:
                c.showPage()
                y = height - 40

        c.save()
        return str(file_path)

    @staticmethod
    async def export_docx(
        db: AsyncSession,
        project_id: int,

    ) -> str:
        review_record = await TrainingContentService.get_latest_phase_output(
            db=db,
            project_id=project_id,
            phase="review",
        )
        if not review_record:
            raise RuntimeError("Review phase output is required for export.")

        project = await ProjectService.get_project(db, project_id)
        if not project:
            raise RuntimeError("Project not found.")

        document = Document()
        document.add_heading(project.title, level=1)
        document.add_paragraph(f"Review Score: {review_record.content.get('score', 'N/A')}")
        document.add_heading("Strengths", level=2)
        for item in review_record.content.get('strengths', []):
            document.add_paragraph(item, style="List Bullet")
        document.add_heading("Improvements", level=2)
        for item in review_record.content.get('improvements', []):
            document.add_paragraph(item, style="List Bullet")
        document.add_paragraph(f"Approved: {review_record.content.get('approved', False)}")

        filename = f"project_{project_id}_review_export_{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.docx"
        file_path = EXPORT_DIR / filename
        document.save(file_path)
        return str(file_path)

    @staticmethod
    async def export_ppt(
        db: AsyncSession,
        project_id: int,
    ) -> str:
        review_record = await TrainingContentService.get_latest_phase_output(
            db=db,
            project_id=project_id,
            phase="review",
        )
        if not review_record:
            raise RuntimeError("Review phase output is required for export.")

        project = await ProjectService.get_project(db, project_id)
        if not project:
            raise RuntimeError("Project not found.")

        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = project.title
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"Review score: {review_record.content.get('score', 'N/A')}"

        strengths_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        strengths_slide.shapes.title.text = "Strengths"
        strengths_body = strengths_slide.shapes.placeholders[1].text_frame
        for item in review_record.content.get('strengths', []):
            strengths_body.add_paragraph().text = item

        improvements_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        improvements_slide.shapes.title.text = "Improvements"
        improvements_body = improvements_slide.shapes.placeholders[1].text_frame
        for item in review_record.content.get('improvements', []):
            improvements_body.add_paragraph().text = item

        filename = f"project_{project_id}_review_export_{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.pptx"
        file_path = EXPORT_DIR / filename
        presentation.save(file_path)
        return str(file_path)
